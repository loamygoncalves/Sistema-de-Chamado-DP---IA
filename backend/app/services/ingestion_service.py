"""Extração e chunking de documentos corporativos para alimentar o RAG.

Suporta PDF, DOCX, XLSX, CSV, PPTX e TXT — inclui os dois formatos lidos
diretamente da pasta de rede local (`local_folder_sync_service.py`): TXT e
PDF. Cada chunk vira um ponto no Qdrant com metadados de rastreabilidade
(documento de origem, página/planilha, departamento).
"""

import csv
import io

import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

CHUNK_SIZE_CHARS = 1500
CHUNK_OVERLAP_CHARS = 200


def _chunk_text(text: str) -> list[str]:
    text = " ".join(text.split())
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + CHUNK_SIZE_CHARS
        chunks.append(text[start:end])
        start = end - CHUNK_OVERLAP_CHARS
    return chunks


def extract_pdf(content: bytes) -> list[dict]:
    reader = PdfReader(io.BytesIO(content))
    chunks = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        for chunk in _chunk_text(text):
            chunks.append({"text": chunk, "metadata": {"page": page_number}})
    return chunks


def extract_docx(content: bytes) -> list[dict]:
    doc = DocxDocument(io.BytesIO(content))
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": chunk, "metadata": {}} for chunk in _chunk_text(full_text)]


def extract_xlsx(content: bytes) -> list[dict]:
    sheets = pd.read_excel(io.BytesIO(content), sheet_name=None, dtype=str)
    chunks = []
    for sheet_name, df in sheets.items():
        rows_text = df.fillna("").astype(str).apply(lambda row: " | ".join(row), axis=1)
        text = "\n".join(rows_text.tolist())
        for chunk in _chunk_text(text):
            chunks.append({"text": chunk, "metadata": {"sheet": sheet_name}})
    return chunks


def extract_csv(content: bytes) -> list[dict]:
    decoded = content.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(decoded))
    rows = [" | ".join(row) for row in reader]
    text = "\n".join(rows)
    return [{"text": chunk, "metadata": {}} for chunk in _chunk_text(text)]


def extract_pptx(content: bytes) -> list[dict]:
    presentation = Presentation(io.BytesIO(content))
    chunks = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            texts.append(slide.notes_slide.notes_text_frame.text)
        slide_text = "\n".join(texts)
        for chunk in _chunk_text(slide_text):
            chunks.append({"text": chunk, "metadata": {"slide": slide_number}})
    return chunks


def extract_txt(content: bytes) -> list[dict]:
    text = content.decode("utf-8", errors="replace")
    return [{"text": chunk, "metadata": {}} for chunk in _chunk_text(text)]


EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "xlsx": extract_xlsx,
    "csv": extract_csv,
    "pptx": extract_pptx,
    "txt": extract_txt,
}


def extract_chunks(file_type: str, content: bytes) -> list[dict]:
    extractor = EXTRACTORS.get(file_type)
    if extractor is None:
        raise ValueError(f"Tipo de documento não suportado: {file_type}")
    return extractor(content)

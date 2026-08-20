import io

import pytest
from pptx import Presentation

from app.services.ingestion_service import extract_pptx


def _build_sample_pptx() -> bytes:
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[1]

    slide1 = presentation.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Banco de Horas"
    slide1.placeholders[1].text_frame.text = "As horas excedentes são compensadas em até 6 meses."

    slide2 = presentation.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Vale Refeição"
    slide2.placeholders[1].text_frame.text = "Cartão definitivo enviado até a 1ª semana do mês seguinte."

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_extract_pptx_returns_chunk_per_slide_with_text():
    chunks = extract_pptx(_build_sample_pptx())

    assert len(chunks) == 2
    assert chunks[0]["metadata"]["slide"] == 1
    assert "Banco de Horas" in chunks[0]["text"]
    assert chunks[1]["metadata"]["slide"] == 2
    assert "Vale Refeição" in chunks[1]["text"]


def test_extract_pptx_skips_empty_slides():
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])  # layout em branco, sem texto
    buffer = io.BytesIO()
    presentation.save(buffer)

    chunks = extract_pptx(buffer.getvalue())
    assert chunks == []


@pytest.mark.parametrize("file_type", ["pptx"])
def test_extract_chunks_dispatches_to_pptx(file_type):
    from app.services.ingestion_service import extract_chunks

    chunks = extract_chunks(file_type, _build_sample_pptx())
    assert len(chunks) == 2
